import os
import sys
import subprocess

def install_pptx():
    print("[PPT Generator] Checking python-pptx installation...")
    try:
        import pptx
    except ImportError:
        print("[PPT Generator] Installing python-pptx...")
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        except Exception as e:
            print(f"[PPT Generator] Error installing package: {e}")
            sys.exit(1)

def build_presentation():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    
    print("[PPT Generator] Constructing slide deck structure...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Theme settings (glowing glassmorphism dark template)
    BG_COLOR = RGBColor(12, 11, 20)        # Dark Indigo
    TITLE_COLOR = RGBColor(243, 244, 246)   # White
    ACCENT_COLOR = RGBColor(167, 139, 250) # Light Purple / Violet
    TEXT_COLOR = RGBColor(156, 163, 175)   # Light Gray

    def set_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_slide_title(slide, text):
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Calibri"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        return title_box

    def add_split_slide(title, bullets, image_filename=None):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        set_bg(slide)
        add_slide_title(slide, title)
        
        # Text block (left half)
        text_width = Inches(5.8) if image_filename else Inches(11.7)
        text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), text_width, Inches(5.3))
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        for idx, (b_title, b_desc) in enumerate(bullets):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            p.text = f"• {b_title}: "
            p.font.name = "Calibri"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = TITLE_COLOR
            p.space_after = Pt(8)
            
            run = p.add_run()
            run.text = b_desc
            run.font.bold = False
            run.font.size = Pt(15)
            run.font.color.rgb = TEXT_COLOR
            p.line_spacing = 1.15
            
        # Image block (right half)
        if image_filename:
            image_path = os.path.join(os.path.dirname(__file__), image_filename)
            if os.path.exists(image_path):
                try:
                    # Let slide picture auto-calculate height keeping aspect ratio
                    slide.shapes.add_picture(image_path, Inches(7.0), Inches(1.5), width=Inches(5.5))
                except Exception as e:
                    print(f"[PPT Generator] Warning: Couldn't insert {image_filename}: {e}")
            else:
                # Add placeholder warning box if image doesn't exist
                placeholder = slide.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.0))
                ptf = placeholder.text_frame
                ptf.word_wrap = True
                pp = ptf.paragraphs[0]
                pp.text = f"[Visual Asset Missing: {image_filename}]"
                pp.font.color.rgb = RGBColor(239, 68, 68)
                pp.font.bold = True

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    
    t_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.2))
    tf = t_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = "SENTRYEYE"
    p.font.name = "Impact"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    p2 = tf.add_paragraph()
    p2.text = "ViT-Based Deepfake Image Detection with Explainable AI & Secure Authentication"
    p2.font.name = "Calibri"
    p2.font.size = Pt(22)
    p2.font.color.rgb = ACCENT_COLOR
    p2.space_before = Pt(8)
    
    m_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.2))
    tf_meta = m_box.text_frame
    tf_meta.word_wrap = True
    tf_meta.margin_left = tf_meta.margin_top = tf_meta.margin_right = tf_meta.margin_bottom = 0
    
    p_sub = tf_meta.paragraphs[0]
    p_sub.text = "Submitted By:\n• Vikash Tiwari (12500121106)\n• Raj Agrahari (12500221053)\n• Abhishek Kumar Mishra (12500121051)"
    p_sub.font.name = "Calibri"
    p_sub.font.size = Pt(15)
    p_sub.font.color.rgb = TEXT_COLOR
    p_sub.line_spacing = 1.2
    
    p_guide = tf_meta.add_paragraph()
    p_guide.text = "Under the Guidance of:\nMr. Sajal Chakraborty, Assistant Professor, Department of CSE\nBengal College of Engineering and Technology, Durgapur"
    p_guide.font.name = "Calibri"
    p_guide.font.size = Pt(15)
    p_guide.font.color.rgb = TEXT_COLOR
    p_guide.space_before = Pt(12)
    p_guide.line_spacing = 1.2

    # --- Slide 2: Introduction ---
    bullets_s2 = [
        ("The Deepfake Challenge", "Generative adversarial networks (GANs) and advanced diffusion models allow synthesis of highly realistic fake human profile pictures and biometric markers."),
        ("Biometric Vulnerability", "Manipulated profiles easily spoof standard face scanners, bypass remote KYC compliance systems, and spread misinformation across digital channels."),
        ("Local Texture Limitations", "Older CNN structures examine local pixels and fail to identify discrepancies in global symmetries, such as pupil alignments or lighting balances."),
        ("The SentryEye Approach", "Bridges this vulnerability by applying Vision Transformers (ViT) to check global structures, backed by Explainable AI (XAI) overlays to visualize detection clues.")
    ]
    add_split_slide("Introduction & Background", bullets_s2)

    # --- Slide 3: Scope and Objectives ---
    bullets_s3 = [
        ("High-Accuracy Classification", "Provide a fine-tuned deep learning network that achieves over 90% detection accuracy on highly compressed face images."),
        ("Visual Interpretability (XAI)", "Implement self-attention rollout mapping to highlight facial pixels containing inconsistencies, eliminating black-box models."),
        ("Administrative Security", "Establish role-based login and signup screens with SHA256 password hashing and Bearer token check integrations to protect APIs."),
        ("Widescreen Usability", "Create a dashboard containing drag-and-drop actions, scan history timelines, and diagnostic gauges for easy presentation.")
    ]
    add_split_slide("Project Objectives & Scope", bullets_s3)

    # --- Slide 4: System Architecture (DFD) ---
    bullets_s4 = [
        ("Data Routing Design", "Coordinates incoming multipart file uploads via FastAPI handlers, sending raw streams directly to preprocessing blocks."),
        ("FastAPI Backend", "Runs as an asynchronous Python server, protecting routes with Bearer token authentication logic."),
        ("Web UI Client", "Constructed using responsive glassmorphic layouts, storing session tokens securely inside localStorage."),
        ("Visualization pipeline", "Decodes self-attention arrays back into RGB color scales and returns base64 string maps.")
    ]
    add_split_slide("SentryEye Architecture & DFD", bullets_s4, "dfd_diagram.jpg")

    # --- Slide 5: System Processing Flow ---
    bullets_s5 = [
        ("Payload Validation", "Intercepts files on drop, confirming that they are valid JPEG/PNG formats prior to backend routing."),
        ("Async FastAPI Handler", "Decodes the raw stream, caches the tensors, and calls model inference dependencies in a thread-pool."),
        ("ViT Inference Engine", "Computes logit probabilities to determine Real/Fake labels and extracts final-block attention matrices."),
        ("Colormap Blending", "Converts attention grids into colormapped heat overlays, generating base64 maps for the dashboard.")
    ]
    add_split_slide("System Processing Flowchart", bullets_s5, "flowchart_diagram.jpg")

    # --- Slide 6: ViT Pipeline ---
    bullets_s6 = [
        ("16x16 Patch Slicing", "Converts the 224x224 RGB input image into a sequence of 196 flat spatial patches ($N=196$)."),
        ("Linear Projection", "Projects flat patches into a 768-dimensional token space matching the encoder's internal dimensions."),
        ("Learnable CLS Token", "Prepends a special [CLS] token at index 0. This token aggregates information across all layers to determine classification."),
        ("Multi-Head Attention", "Utilizes 12 self-attention heads in each encoder block to map global contextual dependencies concurrently.")
    ]
    add_split_slide("Vision Transformer (ViT) Pipeline", bullets_s6, "vit_architecture.jpg")

    # --- Slide 7: Explainable AI Mapping ---
    bullets_s7 = [
        ("Attention Slicing", "Extracts attention weights from the final layer of the transformer block, focusing on relationships from the [CLS] token to the other 196 patch tokens."),
        ("Attention Aggregation", "Averages the attention weights across all 12 heads to produce a consolidated 1D vector of size 196."),
        ("Spatial Bilinear Upscaling", "Reshapes the 196 values to a 14x14 grid and scales it back to 224x224 using bilinear interpolation."),
        ("Transparency Blending", "Superimposes the heatmap overlay using a Jet colormap at 45% opacity, highlighting manipulation hotspots in red.")
    ]
    add_split_slide("Explainable AI (XAI) Mapping", bullets_s7)

    # --- Slide 8: Heatmap Demonstration ---
    bullets_s8 = [
        ("Facial Landmark Focus", "The attention heatmaps demonstrate that SentryEye focuses on key facial boundaries, eyes, and mouth contours."),
        ("Blending Seam Spotting", "Detects subtle skin lighting anomalies and mismatched textures around the eyes and nose boundaries."),
        ("Explainability Verification", "Allows review panels and security officers to instantly see which features led to a fake classification."),
        ("Real vs Fake comparison", "Authentic images show uniform low attention, whereas fake images focus heavily on boundaries.")
    ]
    add_split_slide("XAI Heatmap Output Sample", bullets_s8, "xai_heatmap.jpg")

    # --- Slide 9: Authentication & Registry ---
    bullets_s9 = [
        ("Local Persistence", "Stores user profiles in `users.json` within the root folder, avoiding complex database setups."),
        ("SHA-256 Encryption", "Protects passwords by hashing them with SHA-256 before saving to the database file."),
        ("Route Handlers", "Injects Bearer Token checks into FastAPI routes (`predict` and `history`), blocking unauthorized API queries."),
        ("Transitional UI", "Features a dynamic login card that expands to show signup inputs when switching modes.")
    ]
    add_split_slide("User Authentication & Signup", bullets_s9)

    # --- Slide 10: Performance Benchmarks ---
    bullets_s10 = [
        ("FF++ Test Benchmark", "Validation checks on the FaceForensics++ database yield an accuracy of **94.8%** with a 0.941 F1-score."),
        ("Celeb-DF Benchmark", "Testing on Celeb-DF dataset achieves a validation accuracy of **93.1%**."),
        ("Low-Latency Inference", "Average response latency is **~140ms on standard CPUs** and **~15ms on NVIDIA GPUs**."),
        ("Robust Compression", "Maintains strong detection capability on compressed social media files compared to standard CNNs.")
    ]
    add_split_slide("Experimental Performance Benchmarks", bullets_s10)

    # --- Slide 11: Future Scope & Conclusion ---
    bullets_s11 = [
        ("Temporal Video Parsing", "Extend the pipeline to support video streams by extracting temporal frames and processing sequences with LSTMs."),
        ("Multimodal Fusion", "Incorporate synthetic voice clone detectors alongside video scans for a unified defense system."),
        ("Mobile Edge Quantization", "Quantize model weights (FP32 to INT8) to run the detector directly on mobile and edge devices."),
        ("Conclusion Summary", "SentryEye provides a fast, secure web utility that successfully classifies deepfakes with visual explainability.")
    ]
    add_split_slide("Conclusion & Future Scope", bullets_s11)

    # Save PPTX
    out_name = "Detailed_SentryEye_Presentation.pptx"
    prs.save(out_name)
    print(f"[PPT Generator] Slide deck saved successfully as: {out_name}")

if __name__ == "__main__":
    install_pptx()
    build_presentation()
