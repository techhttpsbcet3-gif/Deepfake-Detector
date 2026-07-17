import sys
import os
import subprocess

def install_pptx():
    print("[PPT Generator] Checking python-pptx installation...")
    try:
        import pptx
        print("[PPT Generator] python-pptx is already installed.")
    except ImportError:
        print("[PPT Generator] Installing python-pptx...")
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
            print("[PPT Generator] Package python-pptx installed successfully.")
        except subprocess.CalledProcessError as e:
            print(f"[PPT Generator] Error installing package: {e}")
            sys.exit(1)

def create_presentation():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    print("[PPT Generator] Generating slide deck...")
    
    prs = Presentation()
    # Set to widescreen 16:9 layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Custom color palette (matching SentryEye dashboard dark mode)
    BG_COLOR = RGBColor(12, 11, 20)        # Dark Indigo
    TITLE_COLOR = RGBColor(243, 244, 246)   # Bright White
    ACCENT_COLOR = RGBColor(167, 139, 250) # Light Purple / Violet
    TEXT_COLOR = RGBColor(156, 163, 175)   # Light Gray
    GREEN_COLOR = RGBColor(16, 185, 129)   # Real Green
    RED_COLOR = RGBColor(239, 68, 68)      # Fake Red
    
    # Helper to set slide background
    def set_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    # Helper to add standard slide title
    def add_slide_title(slide, text):
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Calibri"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        return title_box

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[6] # blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    
    # Title Text Box
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = "SENTRYEYE"
    p.font.name = "Impact"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    p2 = tf.add_paragraph()
    p2.text = "ViT-Based Deepfake Image Detection with Explainable AI"
    p2.font.name = "Calibri"
    p2.font.size = Pt(24)
    p2.font.color.rgb = ACCENT_COLOR
    p2.space_before = Pt(10)
    
    # Metadata Text Box
    meta_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.2))
    tf_meta = meta_box.text_frame
    tf_meta.word_wrap = True
    tf_meta.margin_left = tf_meta.margin_top = tf_meta.margin_right = tf_meta.margin_bottom = 0
    
    p_sub = tf_meta.paragraphs[0]
    p_sub.text = "Submitted By:\n• Vikash Tiwari (12500121106)\n• Raj Agrahari (12500221053)\n• Abhishek Kumar Mishra (12500121051)"
    p_sub.font.name = "Calibri"
    p_sub.font.size = Pt(15)
    p_sub.font.color.rgb = TEXT_COLOR
    p_sub.line_spacing = 1.2
    
    p_guide = tf_meta.add_paragraph()
    p_guide.text = "Under the Guidance of:\nMr. Sajal Chakraborty, Assistant Professor, Department of CSE"
    p_guide.font.name = "Calibri"
    p_guide.font.size = Pt(15)
    p_guide.font.color.rgb = TEXT_COLOR
    p_guide.space_before = Pt(12)
    p_guide.line_spacing = 1.2
    
    # --- Slide 2: Introduction ---
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_slide_title(slide, "Introduction & Problem Statement")
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    bullets = [
        ("The Deepfake Threat", "Generative AI models (GANs, Latent Diffusion) can construct photorealistic human faces, making synthetic media indistinguishable from authentic photographs."),
        ("Security Risks", "Deepfakes pose severe challenges to digital identity verification, security authentication, media integrity, and can propagate targeted misinformation."),
        ("Limitations of Classic CNNs", "Standard Convolutional Neural Networks excel at micro-texture anomalies but fail to capture global symmetry and alignment (like lighting mismatch between eyes)."),
        ("SentryEye Solution", "A high-performance pipeline using Vision Transformers (ViT) to model global features coupled with Self-Attention mapping to visualize manipulation hotspots.")
    ]
    
    for title, desc in bullets:
        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
        p.text = f"• {title}: "
        p.font.name = "Calibri"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.space_after = Pt(8)
        
        # Add description text as inline run
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = TEXT_COLOR

    # --- Slide 3: System Architecture ---
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_slide_title(slide, "Proposed System Architecture")
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    bullets = [
        ("FastAPI Backend", "Python-based asynchronous server that coordinates file stream payloads, schedules tensor preprocessing, and loads PyTorch models efficiently."),
        ("Interactive Web Dashboard", "A premium dashboard built using HTML5, Vanilla CSS3 (featuring glassmorphism layouts), and JavaScript (ES6) with drag-and-drop triggers."),
        ("ViT Classification Engine", "Integrates Google's Vision Transformer (ViT-Base-Patch16-224) to evaluate image patch matrices and classify authentic features."),
        ("Explainable AI (XAI)", "Extracts cross-attention grids directly from the final encoder block to compile and render facial manipulation hotspots.")
    ]
    
    for title, desc in bullets:
        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
        p.text = f"• {title}: "
        p.font.name = "Calibri"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.space_after = Pt(8)
        
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = TEXT_COLOR

    # --- Slide 4: Vision Transformer (ViT) Backbone ---
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_slide_title(slide, "Deep Learning Backbone: Vision Transformer (ViT)")
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    bullets = [
        ("Patch Extraction", "The input 224x224 RGB image is sliced into 196 non-overlapping spatial patches of size 16x16 pixels."),
        ("Linear Projection", "Flat patch matrices are projected into a 768-dimensional token space (embedding vectors)."),
        ("CLS Token", "A learnable [CLS] classification token is prepended to accumulate image-wide representations across multi-attention layers."),
        ("Self-Attention Encoder", "Applies 12 encoder blocks with 12 attention heads. Self-Attention evaluates relationships between all patches concurrently, tracing global structural coherence.")
    ]
    
    for title, desc in bullets:
        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
        p.text = f"• {title}: "
        p.font.name = "Calibri"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.space_after = Pt(8)
        
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = TEXT_COLOR

    # --- Slide 5: Explainable AI ---
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_slide_title(slide, "Explainable AI (XAI) Mapping Pipeline")
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    bullets = [
        ("Attention Extraction", "Extracts the self-attention weights tensor directly from the final transformer block ([CLS] scanner token to patch tokens)."),
        ("Grid Resolution Mapping", "The resulting 196 attention scores are reshaped into a 14x14 grid representing localized parts of the face."),
        ("Bilinear Interpolation", "Upscales the 14x14 grid back to the original image dimensions (e.g. 224x224)."),
        ("Alpha Overlay Blending", "Applies a Jet colormap (red for high attention, blue for low) and blends it with the original face using a 0.45 transparency factor.")
    ]
    
    for title, desc in bullets:
        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
        p.text = f"• {title}: "
        p.font.name = "Calibri"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.space_after = Pt(8)
        
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = TEXT_COLOR

    # --- Slide 6: Results & Performance ---
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_slide_title(slide, "Results & Performance Evaluation")
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    bullets = [
        ("High Detection Accuracy", "Achieves an overall classification accuracy of 92.0% (and up to 97.0% validation scores) on Celeb-DF and FaceForensics++."),
        ("Real-Time Latency", "Processing latency averages ~140ms on standard CPUs (and ~15ms on NVIDIA CUDA GPUs)."),
        ("JPEG Robustness", "Maintains solid detection capabilities even under 85% compression quality, outperforming convolutional networks."),
        ("Detailed Indicators", "Displays real/fake verdict badges and confidence percentages side-by-side with XAI overlays.")
    ]
    
    for title, desc in bullets:
        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
        p.text = f"• {title}: "
        p.font.name = "Calibri"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.space_after = Pt(8)
        
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = TEXT_COLOR

    # --- Slide 7: Future Scope & Conclusion ---
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_slide_title(slide, "Future Scope & Conclusion")
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    bullets = [
        ("Video Stream Decoding", "Extension to video formats by parsing temporal consistencies across sequential frames (using LSTMs/GRUs)."),
        ("Multimodal Fusion", "Merging facial anomaly checking with speech cloned-audio detection for unified media protection."),
        ("Edge Deployments", "Implementing model weight quantization (FP32 to INT8) to optimize and load models on mobile devices."),
        ("Conclusion", "SentryEye provides a reliable, asynchronous web tool that successfully detects modern deepfakes with over 90% accuracy while providing visual explainability.")
    ]
    
    for title, desc in bullets:
        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
        p.text = f"• {title}: "
        p.font.name = "Calibri"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.space_after = Pt(8)
        
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = TEXT_COLOR

    # Save presentation
    output_filename = "SentryEye_Presentation.pptx"
    prs.save(output_filename)
    print(f"[PPT Generator] Presentation saved successfully as: {output_filename}")

if __name__ == "__main__":
    install_pptx()
    create_presentation()
